import os
import json
import tempfile
import uuid
import re
import sys
import unicodedata
import asyncio
import hashlib
import threading
import difflib
import time
from datetime import datetime, timezone
from typing import Annotated, Literal
from types import SimpleNamespace

from dotenv import load_dotenv
load_dotenv()

_ERR_EMPTY_QUESTION_IDS = 'question_ids cannot be empty'

import aiofiles
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from google import genai
from google.genai import errors
from sqlalchemy import cast, Integer

from database import SessionLocal, QuestionRecord, SelectedQuestion, LivePushLog, get_next_ques_number, detect_dominant_language, has_wrong_script_chars, normalize_question_key, get_next_ques_numbers_batch, get_next_order_numbers_batch, TAMIL_CHAR_RE, LATIN_CHAR_RE
from external_sync import push_records_to_live_db

app = FastAPI(title='Question Generator API')
app.add_middleware(CORSMiddleware, allow_origins=['*'], allow_methods=['*'], allow_headers=['*'])

def _configure_tesseract():
    import pytesseract
    env_path = os.getenv('TESSERACT_CMD')
    if env_path:
        pytesseract.pytesseract.tesseract_cmd = env_path
        return
    if sys.platform.startswith('win'):
        default_win_path = 'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'
        if os.path.exists(default_win_path): pytesseract.pytesseract.tesseract_cmd = default_win_path

def _detect_file_type(file_path: str):
    _, extension = os.path.splitext(file_path)
    return extension.lower()

def _extract_text_from_pdf(file_path: str, ocr_lang: str='eng+tam'):
    import fitz
    import pytesseract
    from PIL import Image
    import io
    _configure_tesseract()
    full_text = ''
    doc = fitz.open(file_path)
    for page in doc:
        page_text = page.get_text()
        if page_text.strip():
            full_text += page_text + '\n'
        else:
            pix = page.get_pixmap(dpi=400)
            page_image = Image.open(io.BytesIO(pix.tobytes('png')))
            ocr_text = pytesseract.image_to_string(page_image, lang=ocr_lang)
            if ocr_text.strip(): full_text += ocr_text + '\n'
    doc.close()
    return full_text

def _extract_text_from_docx(file_path: str):
    from docx import Document
    doc = Document(file_path)
    return '\n'.join((p.text for p in doc.paragraphs))

def _extract_text_from_pptx(file_path: str):
    from pptx import Presentation
    prs = Presentation(file_path)
    full_text = ''
    for slide_number, slide in enumerate(prs.slides, start=1):
        full_text += f'\n--- slide{slide_number} --\n'
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs: full_text += run.text + ' '
                full_text += '\n'
    return full_text

def _extract_text_from_txt(file_path: str):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def _extract_text_from_image(file_path: str, ocr_lang: str='eng+tam'):
    import pytesseract
    from PIL import Image
    _configure_tesseract()
    return pytesseract.image_to_string(Image.open(file_path), lang=ocr_lang)

def extract_text(file_path: str, ocr_lang: str='eng+tam'):
    file_type = _detect_file_type(file_path)
    if file_type == '.pdf':
        return _extract_text_from_pdf(file_path, ocr_lang)
    elif file_type == '.docx':
        return _extract_text_from_docx(file_path)
    elif file_type == '.pptx':
        return _extract_text_from_pptx(file_path)
    elif file_type == '.txt':
        return _extract_text_from_txt(file_path)
    elif file_type in ('.jpg', '.jpeg', '.png'):
        return _extract_text_from_image(file_path, ocr_lang)
    else:
        raise ValueError(f"Unsupported file type: '{file_type}'. Supported types: .pdf, .docx, .pptx, .txt, .jpg, .jpeg, .png")

_notes_text_cache: dict[str, str] = {}
_notes_text_cache_lock = threading.Lock()
CHUNK_THRESHOLD = 30
CHUNK_SIZE = 25
# Raised from 5 -> 10: lets more chunk/verify/proofread calls run in parallel
# per request, which is the main lever for cutting wall-clock time on large
# (e.g. 150-question) generations. Override with GEMINI_MAX_CONCURRENCY in
# .env if your Gemini API key's rate limit can't sustain this.
GEMINI_MAX_CONCURRENCY = int(os.getenv('GEMINI_MAX_CONCURRENCY', '10'))
# Questions per generation chunk. Raised from 15 -> 20 so fewer round trips
# are needed for the same total count (e.g. 180 questions = 9 chunks instead
# of 12), without changing prompt content, formats, or validation logic.
GENERATION_CHUNK_SIZE = int(os.getenv('GENERATION_CHUNK_SIZE', '20'))
_MAX_RATE_LIMIT_RETRIES = 5
_RATE_LIMIT_BACKOFF_SECONDS = 5
_RETRY_DELAY_RE = re.compile('retryDelay[\'\\"]?\\s*:\\s*[\'\\"]?(\\d+)')
MAX_NOTES_CHARS = int(os.getenv('MAX_NOTES_CHARS', '60000'))
ALLOWED_GEMINI_MODELS = {'gemini-3.5-flash-lite': 'Gemini 3.5 Flash-Lite - cheapest/fastest Gemini 3 model (default)', 'gemini-3.1-flash-lite': 'Gemini 3.1 Flash-Lite - frontier-class quality, low cost', 'gemini-3.5-flash': 'Gemini 3.5 Flash - most intelligent Flash model, sustained/agentic tasks', 'gemini-3.6-flash': 'Gemini 3.6 Flash - newest, strong agentic/multimodal performance', 'gemini-3-flash-preview': 'Gemini 3 Flash (Preview) - frontier quality at Flash pricing', 'gemini-3.1-pro-preview': 'Gemini 3.1 Pro (Preview) - flagship reasoning model', 'gemini-2.5-flash-lite': 'Gemini 2.5 Flash-Lite - older gen, may 404 on new API keys', 'gemini-2.5-flash': 'Gemini 2.5 Flash - older gen, may 404 on new API keys', 'gemini-2.5-pro': 'Gemini 2.5 Pro - older gen, may 404 on new API keys'}
_FALLBACK_GEMINI_MODEL = 'gemini-3.5-flash-lite'
DEFAULT_GEMINI_MODEL = os.getenv('GEMINI_MODEL', _FALLBACK_GEMINI_MODEL)
if DEFAULT_GEMINI_MODEL not in ALLOWED_GEMINI_MODELS:
    print(f'[startup] GEMINI_MODEL={DEFAULT_GEMINI_MODEL!r} is not in ALLOWED_GEMINI_MODELS; falling back to {_FALLBACK_GEMINI_MODEL!r}.')
    DEFAULT_GEMINI_MODEL = _FALLBACK_GEMINI_MODEL

_COMBINING_MARKS = '\u0B82\u0BBE-\u0BCD\u0BD7\u0900-\u0903\u093A-\u094F\u0951-\u0957\u0962-\u0963'
_DUPLICATED_MARK_RE = re.compile(f'([{_COMBINING_MARKS}])\\1+')
_TAMIL_CHAR_RE = re.compile('[\\u0B80-\\u0BFF]')
_LATIN_CHAR_RE = re.compile('[A-Za-z]')

_WORD_NUMBER_LABEL_RE = re.compile(r'^\s*(plus|minus)\s+\S+\s*$', re.IGNORECASE)
_BAR_WRAP_STOPWORDS = r'(?:a|an|the|its|this|that|all|any|each|both|either|neither|of|is|are|was|were)'

def clean_math_formatting(text: str) -> str:
    if not text: return text
    if _WORD_NUMBER_LABEL_RE.match(text):
        return text
    
    text = re.sub(r'\babsolute\s+value\s+of\s+([A-Za-z0-9_]+)\s+minus\s+([A-Za-z0-9_]+)\s+equals\s+([A-Za-z0-9_]+)\b', r'|\1 - \2| = \3', text, flags=re.IGNORECASE)
    text = re.sub(r'\babsolute\s+value\s+of\s+(?!' + _BAR_WRAP_STOPWORDS + r'\b)([A-Za-z0-9_ -]+?)\s+equals\s+([A-Za-z0-9_ -]+?)\b', r'|\1| = \2', text, flags=re.IGNORECASE)
    text = re.sub(r'\babsolute\s+value\s+of\s+(?!' + _BAR_WRAP_STOPWORDS + r'\b)([A-Za-z0-9_][A-Za-z0-9_\^-]*)', r'|\1|', text, flags=re.IGNORECASE)
    text = re.sub(r'\bmodulus\s+of\s+(?!' + _BAR_WRAP_STOPWORDS + r'\b)([A-Za-z0-9_][A-Za-z0-9_\^-]*)', r'|\1|', text, flags=re.IGNORECASE)
    
    def clean_division(match):
        op1 = match.group(1).strip()
        op2 = match.group(2).strip()
        if '+' in op1 or '-' in op1:
            op1 = f"({op1})"
        if '+' in op2 or '-' in op2:
            op2 = f"({op2})"
        return f"{op1} / {op2}"
        
    text = re.sub(r'(\b[A-Za-z0-9_+-]+)\s+divided\s+by\s+([A-Za-z0-9_+-]+\b)', clean_division, text, flags=re.IGNORECASE)
    text = re.sub(r'(\b[A-Za-z0-9_+-]+)\s+multiplied\s+by\s+([A-Za-z0-9_+-]+\b)', r'\1 × \2', text, flags=re.IGNORECASE)
    text = re.sub(r'(\b[A-Za-z0-9_]+?)\s+minus\s+(\b[A-Za-z0-9_]+?\b)', r'\1 - \2', text, flags=re.IGNORECASE)
    text = re.sub(r'(\b[A-Za-z0-9_]+?)\s+plus\s+(\b[A-Za-z0-9_]+?\b)', r'\1 + \2', text, flags=re.IGNORECASE)
    text = re.sub(r'\bminus\b', '-', text, flags=re.IGNORECASE)
    text = re.sub(r'\bplus\b', '+', text, flags=re.IGNORECASE)
    text = re.sub(r'\bis\s+not\s+equal\s+to\b', '≠', text, flags=re.IGNORECASE)
    text = re.sub(r'\bnot\s+equal\s+to\b', '≠', text, flags=re.IGNORECASE)
    text = re.sub(r'\bis\s+equal\s+to\b', '=', text, flags=re.IGNORECASE)
    text = re.sub(r'\bequals\b', '=', text, flags=re.IGNORECASE)
    text = re.sub(r'\b0\s+to\s+2p\b', '0 to 2π', text, flags=re.IGNORECASE)
    text = re.sub(r'-p\s+to\s+p\b', '-π to π', text, flags=re.IGNORECASE)
    text = re.sub(r'-p/2\s+to\s+p/2\b', '-π/2 to π/2', text, flags=re.IGNORECASE)
    text = re.sub(r'\b([zxyab])([1-9])\b', r'\1_\2', text, flags=re.IGNORECASE)
    text = re.sub(r'\bz\s+bar\b', 'z̅', text, flags=re.IGNORECASE)
    text = re.sub(r'\bzbar\b', 'z̅', text, flags=re.IGNORECASE)
    text = re.sub(r'\b([A-Za-z])_bar\b', lambda m: m.group(1) + '\u0305', text, flags=re.IGNORECASE)

    latex_greek_map = {
        r'\alpha': 'α', r'\beta': 'β', r'\gamma': 'γ', r'\theta': 'θ',
        r'\lambda': 'λ', r'\mu': 'μ', r'\pi': 'π', r'\phi': 'φ',
        r'\sigma': 'σ', r'\omega': 'ω', r'\delta': 'δ', r'\epsilon': 'ε',
        r'\eta': 'η', r'\psi': 'ψ', r'\tau': 'τ', r'\chi': 'χ',
        r'\xi': 'ξ', r'\zeta': 'ζ'
    }
    for lat, sym in latex_greek_map.items():
        text = re.sub(re.escape(lat), sym, text, flags=re.IGNORECASE)
    
    word_greek_map = {
        'alpha': 'α', 'beta': 'β', 'gamma': 'γ', 'theta': 'θ',
        'lambda': 'λ', 'mu': 'μ', 'pi': 'π', 'phi': 'φ',
        'sigma': 'σ', 'omega': 'ω', 'delta': 'δ', 'epsilon': 'ε',
        'eta': 'η', 'psi': 'ψ', 'tau': 'τ', 'chi': 'χ',
        'xi': 'ξ', 'zeta': 'ζ'
    }
    for word, sym in word_greek_map.items():
        text = re.sub(r'\b' + re.escape(word) + r'\b', sym, text, flags=re.IGNORECASE)
    
    text = text.replace(r'\times', '×')
    text = text.replace(r'\cdot', '·')
    text = text.replace(r'\pm', '±')
    text = text.replace(r'\neq', '≠')
    text = text.replace(r'\le', '≤')
    text = text.replace(r'\ge', '≥')
    text = text.replace(r'\infty', '∞')
    text = text.replace(r'\approx', '≈')
    text = text.replace(r'\subseteq', '⊆')
    text = text.replace(r'\cap', '∩')
    text = text.replace(r'\cup', '∪')
    text = text.replace(r'\in', '∈')
    text = re.sub(r'\\bar\{([^{}]+)\}', r'bar(\1)', text)
    text = re.sub(r'\\bar\s+([A-Za-z0-9_]+)', r'bar(\1)', text)
    text = re.sub(r'\\vec\{([^{}]+)\}', r'vec(\1)', text)
    text = re.sub(r'\\hat\{([^{}]+)\}', r'hat(\1)', text)
    text = re.sub(r'\$(.*?)\$', r'\1', text)
    text = text.replace('$$', '')
    text = re.sub(r'\^{(.*?)\}', r'^\1', text)
    text = re.sub(r'_\{(.*?)\}', r'_\1', text)
    text = re.sub(r'(\b[A-Za-z0-9_]+(?:\^[-0-9T]+)?)\s*\*\s*([A-Za-z0-9_]+(?:\^[-0-9T]+)?\b)', r'\1 \2', text)
    text = re.sub(r'(\)\^[-0-9T]+)\s*\*\s*([A-Za-z0-9_]+)', r'\1 \2', text)
    text = re.sub(r'([A-Za-z0-9_]+)\s*\*\s*(\(\b)', r'\1 \2', text)
    text = re.sub(r'\bz\s*\*\s*bar\(z\)', r'z bar(z)', text)
    text = re.sub(r'\bz\s*\*\s*z\s+bar', r'z z bar', text)
    
    return text

def clean_extraction_artifacts(text: str):
    if not text: return text
    text = unicodedata.normalize('NFC', text)
    text = _DUPLICATED_MARK_RE.sub('\\1', text)
    return text

def clean_question_payload(data: dict) -> dict:
    if not data or not isinstance(data, dict):
        return data
    questions = data.get('questions', [])
    if not isinstance(questions, list):
        return data
    for q in questions:
        if not isinstance(q, dict):
            continue
        if 'question' in q and isinstance(q['question'], str):
            q['question'] = clean_math_formatting(clean_extraction_artifacts(q['question']))
        if 'options' in q and isinstance(q['options'], dict):
            for opt_key, opt_val in q['options'].items():
                if isinstance(opt_val, str):
                    q['options'][opt_key] = clean_math_formatting(clean_extraction_artifacts(opt_val))
        if 'explanation' in q and isinstance(q['explanation'], str):
            q['explanation'] = clean_math_formatting(clean_extraction_artifacts(q['explanation']))
        if 'correct_answer_text' in q and isinstance(q['correct_answer_text'], str):
            q['correct_answer_text'] = clean_math_formatting(clean_extraction_artifacts(q['correct_answer_text']))
    return data

def _norm_compare(text: str) -> str:
    return re.sub('\\s+', ' ', (text or '').strip().lower())

def fix_answer_consistency(data: dict):
    for q in data.get('questions', []):
        correct_text = q.get('correct_answer_text')
        options = q.get('options', {})
        current_letter = q.get('answer')
        if not correct_text or current_letter not in options: continue
        if _norm_compare(options.get(current_letter, '')) != _norm_compare(correct_text):
            for letter, text in options.items():
                if _norm_compare(text) == _norm_compare(correct_text):
                    q['answer'] = letter
                    break
    return data

def build_prompt(notes_text: str, count: int, difficulty: str, language: str, formats_list: list[int], avoid_questions: list[str] | None=None):
    difficulty_guides = {
        'easy': "Direct factual recall. One correct fact, no reasoning needed. (Matches TNPSC Group 4 / TNTET difficulty.)",
        'moderate': "Requires connecting two related facts from the notes, or a 'which of the following is/is not correct' style statement question. (Matches TNPSC Group 2 difficulty.)",
        'hard': "Requires analysis, comparison, or applying a concept to a new situation. Multi-statement or assertion-reason format. (Matches UPSC Prelims/Mains difficulty.)"
    }
    avoid_block = ''
    if avoid_questions:
        numbered = '\n'.join((f'- {q}' for q in avoid_questions))
        avoid_block = f'\nIMPORTANT - DO NOT REPEAT THESE (already generated, in use elsewhere):\nThe following questions already exist for this same notes/difficulty\ncombination. Do NOT generate anything that tests the same underlying fact\nas any of these, even if reworded differently. Every question you write\nnow must cover a fact or angle NOT already covered below.\n\nALREADY-USED QUESTIONS:\n{numbered}\n\n'

    FORMAT_NAMES = {
        1: "Choose the Correct Answer (MCQ): Standard direct multiple-choice question.",
        2: "Fill in the Blanks: Question stem contains an underscore ______ blank, and options A, B, C, D contain possible words to fill the blank. Do not write the word 'dash' or 'blank', use the underscore symbol ______ itself.",
        3: "Match the Following: Question stem lists EXACTLY 4 numbered items on the left and EXACTLY 4 lettered items on the right (1-4 and a-d, or i-iv, matching whatever numbering the source uses) - never 3, never 5, always 4 on each side. The 4 options A, B, C, D each show one complete, different matching combination (e.g., '1-c, 2-a, 3-d, 4-b'), and every option must pair up all 4 items, not a subset.",
        4: "Complete the Sentence: Question stem starts a sentence that requires completion, options contain the concluding parts of the sentence.",
        5: "Identify the Correct Statement: Question lists EXACTLY 2 short numbered statements (1. and 2.) and asks to identify which is/are correct. Options are e.g., '1 only', '2 only', 'Both 1 and 2', 'None' (in Tamil: '1 மட்டும் சரி', '2 மட்டும் சரி', '1 மற்றும் 2 மட்டும் சரி', etc.). Keep it to 2 statements, not 3 - stays quick to read.",
        6: "Name the Following: Question stem describes a process, concept, or entity, and options list possible names/labels.",
        7: "Which of the following: Question stem starts with 'Which of the following...' (e.g., 'Which of the following is/is not true regarding...?').",
        8: "Assertion and Reason: Question stem has 'Assertion (A):' and 'Reason (R):' (in Tamil: 'கூற்று (A):' and 'காரணம் (R):'). Options must strictly follow standard templates.",
        9: "Pick The Odd One Out: Question stem presents 4 terms/statements, and asks the candidate to identify the odd one out.",
        10: "Correct the Incorrect Statement: Question stem states an incorrect fact/statement, and options offer corrected versions of that statement.",
        11: "Where-based questions: Question stem begins with 'Where...' (in Tamil: 'எங்கு...').",
        12: "Who-based questions: Question stem begins with 'Who...' (in Tamil: 'யார்...').",
        13: "What-based questions: Question stem begins with 'What...' (in Tamil: 'என்ன...').",
        14: "Which-based questions: Question stem begins with 'Which...' (in Tamil: 'எது...' / 'எவை...')."
    }

    formatting_lines = []
    for idx, fmt_num in enumerate(formats_list, start=1):
        formatting_lines.append(f"Question {idx} (index {idx-1} in the JSON 'questions' array) MUST be of Format {fmt_num}: {FORMAT_NAMES.get(fmt_num, 'MCQ')}")
    formatting_instructions_block = "\n".join(formatting_lines)

    format_section = f"""
IMPORTANT - QUESTION FORMAT DIVERSITY RULES:
You must distribute the generated questions across a diverse, balanced mix of the 14 distinct question formats. Do not rely heavily on any single format; strive for a balanced distribution of these types throughout the generated questions (where supported by the notes).

You must distribute the generated questions across a diverse mix of the 14 formats listed above. The formats should NOT be restricted by the selected difficulty level; rather, all 14 formats are fully active and must be used across all difficulties (easy, moderate, hard). The selected difficulty level '{difficulty}' only determines the conceptual complexity and depth of the question content, as defined in the difficulty guide.
"""

    prompt = f'''You are a veteran question-setter for Indian government exams: TNPSC (Group 1/2/4), UPSC (Prelims/Mains), and TNTET. You are creating a practice quiz from study notes, matching real exam patterns exactly for the requested difficulty band.
{avoid_block}

IMPORTANT - LANGUAGE RULE (already determined - do not re-decide this):
The STUDY NOTES below are written PRIMARILY in {language}. Write the ENTIRE quiz - every question, all four options, and any explanations - in {language} ONLY, using its native script. Do not translate anything into a different language, and do not switch language partway through.

The STUDY NOTES themselves may contain a few individual words or phrases in a different script (for example a proper noun, place name, title, or quoted term) - seeing that in the SOURCE does NOT mean you should switch the response language. But this works only one way: your OUTPUT must not contain ANY other-script word or phrase either, for ANY reason. Concretely, this means:
- Do NOT add an English word in parentheses after a {language} term to "help the reader understand" (e.g. do not write the {language} word followed by its English gloss in brackets). The learner reading this already reads {language}; a parenthetical gloss is not requested and is not allowed.
- Do NOT keep a technical/scientific term in English/Latin script inside an otherwise-{language} sentence. Transliterate or translate it into {language} script like the rest of the sentence.
- The ONLY exceptions are: numerals (0-9), standard mathematical/scientific symbols, and a proper noun that has no established {language}-script form at all (rare) - and even then, write it using {language} script transliteration if any conventional transliteration exists.
Every single question and every single option you write must be entirely in {language} script, with no exceptions. If {language} is Tamil, this means not one single English word, English letter, or English abbreviation may appear anywhere in the question, options, or explanation - not even a unit, a technical term, or a single letter used as a label.

IMPORTANT - TEXT QUALITY RULE:
The STUDY NOTES below were extracted automatically from a document and may contain minor extraction artifacts (jumbled vowel signs, jumbled or kissing characters, jumbled script diacritics, jumbled spacing, jumbled script diacritics, or similar glitches) - this is a known limitation of automated text extraction, especially for complex scripts. Do NOT copy such errors into your output, even when you are quoting an exact line, poem, or verbatim passage from the notes - correct the spelling in quoted material too. Use correct, standard spelling and grammar for the detected language in every question, option, and explanation, based on your own knowledge of proper spelling - even if the source text contains extraction noise. Pay special attention to: (a) compound words (two words joined together), which are especially prone to losing a letter at the joint, and (b) grammatical suffix endings (verb/noun inflections), which can also lose a consonant when extracted (for example "இருந்து" appearing where "இருந்தது" is grammatically required, or "நெடிலாக்" where "நெடிலாகக்" is required). After writing each sentence, mentally check that every word is a complete, correctly spelled word in the detected language - not a fragment. Base the FACTS strictly on the notes, but express them in clean, correctly spelled language.

IMPORTANT - MATHEMATICAL FORMATTING RULES:
1. For math questions, NEVER wrap variables, expressions, or equations in dollar symbols (do NOT use $A$ or $$A$$ or $x^2$). Instead, write them simply as plain, standard text (e.g., write A, x^2, y_1, I_n, AA^T = A^TA = I_n).
2. ALWAYS use standard mathematical operators (+, -, ×, /, =, <, >, ≤, ≥, √, ^) and symbols instead of writing them out in words. Do NOT write English words for operations or equations:
   - NEVER write 'multiplied by' or 'times'; use × or simple space/juxtaposition (e.g., write 'z z_bar' or 'z × z_bar', not 'z multiplied by z bar'; write '2 × 3', not '2 multiplied by 3').
   - NEVER write 'divided by' or 'division'; use the division slash / (e.g., write '(3 + 4i) / (5 - 12i)', not '3+4i divided by 5-12i').
   - NEVER write 'minus'; use the minus sign - (e.g., write 'a - b', not 'a minus b').
   - NEVER write 'plus'; use the plus sign + (e.g., write 'a + b', not 'a plus b').
   - NEVER write 'equals' or 'is equal to'; use the equals sign = (e.g., write '|z - z_0| = r', not 'absolute value of z minus z0 is equal to r').
   - NEVER write 'absolute value of' or 'modulus of'; use vertical bars |...| (e.g., write '|z| = 1', not 'modulus of z is equal to 1').
   - NEVER write 'square root of'; use the radical symbol √ or power ^(1/2) (e.g., write '√2', not 'square root of 2').
3. NEVER use LaTeX style backslash commands (do NOT write \\alpha, \\beta, \\theta, \\lambda, \\times, \\cdot).
4. ALWAYS use clean mathematical Unicode symbols for Greek letters (e.g., write α, β, θ, λ, ω, π) and operators. NEVER write them out in English words like "alpha", "beta", "theta", "omega", "pi", or single letters like "p" (always use the Unicode character π, not p or pi).
5. Use proper subscripts like z_1, z_2, a_ij instead of writing them as z1, z2, aij.
6. Make sure math equations look clean, natural, and readable as if written on a plain paper test sheet by a human examiner. No raw code snippets, no markup symbols.

IMPORTANT - NEVER REFER TO THE SOURCE MATERIAL:
The notes below are your private research, not something the learner has read or will ever see. Never let a question, option, or explanation refer to the notes/passage/text/document/material/source/chapter/lesson/section/page in any way - no "according to the notes", "as mentioned in the passage", "as per the given text", "based on the above", "in the document", "as discussed", "the notes state", "as given above", "from the notes", "in this lesson", "in this chapter", or any equivalent phrase in {language} or any other language. This applies to EVERY field, including the explanation field. Write each question and explanation as a plain statement of fact or a plain question about the world - exactly as an experienced examiner would write it from memory, with zero trace that it was derived from a supplied document. If you catch yourself about to write a phrase that references where the fact came from, delete it and restate the sentence as a direct fact or question instead.

This also applies to numbered worked examples inside the notes (common in accounts/maths/science textbooks, e.g. "Illustration 16", "Example 3", "Problem 5", "Case Study 2", "Exercise 4.2"). NEVER write a question like "In Illustration 16, what is the value of X?" - the learner has no way to look up "Illustration 16" and cannot answer it; this makes the question unanswerable and worthless. Instead, pull the actual given data (the specific numbers, names, or scenario details from that worked example) directly into the question stem itself, so the question is fully self-contained - e.g. instead of "In Illustration 16, find the depreciation amount", write "A machine costing ₹50,000 is depreciated at 10% per annum using the straight-line method. What is the annual depreciation?" using the real figures from that example. If a worked example's data is too long to restate within the 18-second limit, pick a shorter fact from the same example instead of citing its label.

Concrete phrases you must NEVER produce, in any field, in any language (this list is illustrative, not exhaustive - the same idea in any wording or any language is banned):
- English: "according to the notes/passage/text/document/above", "as mentioned/stated/given/discussed in the notes/passage/above", "from the notes/passage/document", "in this lesson/chapter/section/page/passage", "the notes state/say/mention", "based on the given text/above", "as per the above"
- English (worked-example labels): "in Illustration 16", "as per Example 3", "in Problem 5", "in Case Study 2", "as solved in Exercise 4.2", or any other bare reference to a numbered illustration/example/problem/case/exercise without restating its actual data
- Tamil: "இந்த பக்கத்தில்", "பக்கத்தில் இருந்து", "இந்த குறிப்பில்", "குறிப்புகளின்படி" / "குறிப்புகளின் படி", "இந்த பாடத்தில்", "மேலே கொடுக்கப்பட்ட", "கொடுக்கப்பட்ட பத்தியில்", "இந்த பத்தியில்"
A well-formed question/explanation reads exactly like a standalone exam question - the learner should have no way of knowing it was generated from a document at all.

IMPORTANT - FORMAT ASSIGNMENT FOR EACH QUESTION:
You must generate exactly {count} questions, where each question in the output JSON array must strictly follow the format assigned to its index:
{formatting_instructions_block}

{format_section}
Only use these formats where they genuinely fit the fact being tested and the notes support it - do not force an assertion-reason or statement format onto a fact that is naturally a simple direct-answer question. Most questions should still be plain direct-answer MCQs.

IMPORTANT - LENGTH AND READABILITY (applies at EVERY difficulty level, especially moderate and hard):
The candidate gets a strict 18 SECONDS on screen to read the question stem, read all four options, and pick an answer, before the screen moves on. These questions are also frequently copied out as plain text and read on a phone screen with no app formatting around them at all - so the raw wording itself, not just how an app displays it, has to be short and simple enough to take in at a glance. Every question must be fully readable and answerable within that 18-second window, on a small screen, with zero re-reading. A moderate or hard question is NEVER made harder by being longer or more elaborately worded - it is made harder by WHAT it asks (connecting two facts, spotting the one incorrect statement, applying a concept), while staying just as quick to read as an easy question. Long, essay-like question stems and long, full-sentence options are a defect, not a sign of difficulty - fix them, do not write them.
Concrete limits (hard caps, not targets to approach):
- Question stem: at most ONE short, simple sentence, roughly 12-15 words maximum. Use plain, direct wording - no throat-clearing phrases like "with respect to", "in relation to", "in the context of", or extra qualifying clauses. If you need Assertion/Reason or numbered statements, each individual statement must be a very short clause (roughly 6-8 words), not a sentence, and use at most TWO statements, not three.
- Each option: a short phrase, term, name, date, or clause - NOT a sentence. Aim for well under 6 words per option unless the fact itself (e.g. a direct quotation the notes require) genuinely cannot be shortened.
- If your first draft of a question or option is long, rewrite it shorter before finalizing - do not add explanatory clauses, justifications, or extra context into the question or the options. Save any elaboration for the separate "explanation" field, which is read AFTER answering and is not time-limited.
- If a fact is inherently too long to fit a 12-15 word stem plus four sub-6-word options within 18 seconds, simplify what is being asked rather than keeping the full detail - a quick, slightly narrower question beats one the candidate cannot finish reading in time.
- Prefer the plainer, shorter question formats (direct MCQ, fill-in-the-blank, name-the-following, which/who/what/where-based) whenever the notes support them. Only reach for the denser formats (Match the Following, multi-statement, Assertion-Reason) when the fact genuinely needs that structure - never as padding, and always keeping every individual item inside them just as short as the caps above require.

Generate EXACTLY {count} multiple-choice questions at the "{difficulty}" difficulty level.

Difficulty guide for "{difficulty}": {difficulty_guides[difficulty]}

Rules:
1. Every fact you test must be verifiable from the notes provided below - do not invent facts - but write the question itself as a standalone, general-knowledge question with all needed data restated inline. Never name or point back to the notes, or to a numbered illustration/example/problem within them, as the source (see the NEVER REFER TO THE SOURCE MATERIAL rule above).
2. Each question must have exactly 4 options (A, B, C, D).
3. Clearly indicate the correct option.
4. Write in a clear, exam-appropriate tone.
5. NO two questions may test the same fact or be reworded versions of each other. Every question must test a genuinely different underlying fact.
6. Spread the questions across the ENTIRE set of notes provided - do not cluster them all around one paragraph or section while ignoring the rest.
7. Each question should help the learner genuinely understand and remember the material - avoid trivial or trick wording; test real comprehension of the concept, not just memorization of an isolated phrase.
8. For every question, write a 2-3 sentence "explanation" field, in {language}, that teaches the underlying fact so the learner walks away understanding it - not just a restatement of the correct option. Explain WHY that option is correct as a plain fact (never phrased as "the notes say" or "according to..."), and briefly note what makes the other options wrong if that helps the concept stick. Keep the same dry, formal exam tone as the questions - no casual language.
9. Keep the question and all four options as SHORT as the LENGTH AND READABILITY section above requires - this applies at every difficulty level. Never make a question harder by making it longer.
10. Also include a "correct_answer_text" field containing the EXACT text of whichever option (A/B/C/D) you marked as the answer, copied character-for-character from that option - this is a self-check, so it must genuinely match the option you picked, not just the option you think is most defensible.
11. Before finalizing each question, solve it yourself from scratch as if you were a student, using only the question and options - do not just write something plausible-sounding. Double-check any arithmetic, ordering, or logic in your head, and make sure the option you mark as correct is the option that is ACTUALLY correct, not merely the one your first instinct picked. If you are generating a worked example inside an explanation (e.g. listing sample arrangements or numbers), verify that example itself is valid and satisfies every constraint the question stated (such as "no repeated digits") before including it.
12. When a question or option contains a multi-digit number, group its digits using ONE consistent, correct convention throughout - either the Indian system (groups of 2 digits after the first 3 from the right, e.g. 12,34,567) or the international system (groups of 3 digits, e.g. 1,234,567) - matching whichever system the notes themselves use. Never mix the two styles within the same number, never insert stray digits or extra comma groups, and double-check that the grouped number still reads back as the exact same value intended.
13. Every "explanation" must reach a complete, specific conclusion - never trail off with vague phrases like "...and so on", "...continues in this way", or similar hand-waving. State the final value, place, or term explicitly, every time.

STUDY NOTES:
{notes_text}

REMINDER before you write anything: every question, option, and explanation must be in {language}, and NONE of them may reference the notes, passage, text, document, or any other source - write every item as a standalone fact or question, exactly as a human examiner would from memory. Double-check both of these for each item before finalizing your answer.

Respond with ONLY valid JSON, no extra commentary, no markdown fences, in exactly this structure:

{{
  "questions": [
    {{
      "question": "...",
      "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
      "answer": "A",
      "correct_answer_text": "...",
      "explanation": "...",
      "pattern": 3
    }}
  ]
}}'''
    return prompt

_gemini_client_cache: dict[str, "genai.Client"] = {}
_gemini_client_cache_lock = threading.Lock()

def _get_gemini_client(api_key: str):
    # Reuse one client per API key instead of constructing a new one on every
    # single chunk/verify/proofread call. Safe, purely additive change - the
    # underlying request behavior (model, prompt, config) is unchanged.
    client = _gemini_client_cache.get(api_key)
    if client is None:
        with _gemini_client_cache_lock:
            client = _gemini_client_cache.get(api_key)
            if client is None:
                client = genai.Client(api_key=api_key)
                _gemini_client_cache[api_key] = client
    return client

def call_gemini(prompt: str, api_key: str, model: str=DEFAULT_GEMINI_MODEL, thinking_budget: int=0, _is_retry: bool=False, _rate_limit_attempt: int=0):
    from google.genai import errors, types
    client = _get_gemini_client(api_key)
    try:
        if thinking_budget > 0:
            config = types.GenerateContentConfig(thinking_config=types.ThinkingConfig(thinking_budget=thinking_budget), response_mime_type='application/json')
        else:
            config = types.GenerateContentConfig(response_mime_type='application/json')
        response = client.models.generate_content(model=model, contents=prompt, config=config)
        return response.text
    except errors.ClientError as e:
        is_unavailable = getattr(e, 'code', None) == 404 and 'no longer available' in str(e).lower()
        if is_unavailable and (not _is_retry) and (model != _FALLBACK_GEMINI_MODEL):
            print(f"[call_gemini] Model '{model}' is unavailable on this API key ({e}); retrying once with fallback '{_FALLBACK_GEMINI_MODEL}'.")
            return call_gemini(prompt, api_key, model=_FALLBACK_GEMINI_MODEL, thinking_budget=thinking_budget, _is_retry=True)
        is_rate_limited = getattr(e, 'code', None) == 429 or 'resource_exhausted' in str(e).lower()
        if is_rate_limited and _rate_limit_attempt < _MAX_RATE_LIMIT_RETRIES:
            match = _RETRY_DELAY_RE.search(str(e))
            wait_seconds = int(match.group(1)) + 3 if match else _RATE_LIMIT_BACKOFF_SECONDS * (_rate_limit_attempt + 1)
            print(f'[call_gemini] Rate-limited by Gemini ({e}); retrying in {wait_seconds}s (attempt {_rate_limit_attempt + 1}/{_MAX_RATE_LIMIT_RETRIES}).')
            time.sleep(wait_seconds)
            return call_gemini(prompt, api_key, model=model, thinking_budget=thinking_budget, _is_retry=_is_retry, _rate_limit_attempt=_rate_limit_attempt + 1)
        raise

_PROOFREAD_CHUNK_SIZE = 25

def _build_proofread_prompt(questions: list[dict]) -> str:
    payload = [{'question': q.get('question', ''), 'options': q.get('options', {}), 'answer': q.get('answer', ''), 'explanation': q.get('explanation', '')} for q in questions]
    return f'You are a strict proofreader for a language quiz. Below is a JSON array\ncontaining multiple-choice questions.\n\nIMPORTANT - MATHEMATICAL FORMATTING (fix if broken, do not introduce): Never wrap\nvariables or expressions in dollar signs ($x$, $$x$$) or LaTeX backslash commands\n(\\alpha, \\times, \\cdot). Use plain Unicode symbols instead: operators + - \u00d7 / = < > \u2264 \u2265\n\u221a ^, Greek letters \u03b1 \u03b2 \u03b3 \u03b8 \u03bb \u03bc \u03c0 \u03c6 \u03c3 \u03c9 \u03b4 \u03b5 \u03b7 \u03c8 \u03c4 \u03c7 \u03be \u03b6 (never spelled out as words like\n"alpha" or "pi"), and vertical bars for absolute value (|x|). Never spell out\noperations in words ("divided by", "minus", "equals", "absolute value of") - use\nthe symbol. If a question already contains $ signs, LaTeX commands, or\nspelled-out operators/Greek letters, silently correct them to clean symbol form\nas part of your fix.\n\nCarefully re-read every single word in every question, option, and\nexplanation. Fix ONLY genuine spelling/typo mistakes, such as:\n- an extra duplicated letter, syllable, or diacritic (e.g. a repeated\n  vowel sign or virama/pulli mark) stuck onto or before a word\n- a MISSING or DROPPED letter\n- a stray extra character, word fragment, or lone consonant sitting\n  between words\n- a multi-digit number whose comma grouping is broken, inconsistent, or\n  mixes the Indian and international grouping styles - fix the grouping\n  so it correctly represents the same numeric value, never change the\n  value itself\n- the literal word "dash" or its transliteration appearing where a\n  fill-in-the-blank underscore line should be - replace it with ______\n- any other obvious typo\n\nDo NOT change facts, meaning, wording style, which option is marked\ncorrect, or the overall structure. Do NOT translate anything. Do NOT add,\nremove, or reorder questions - return exactly the same number of items in\nthe same order.\n\nReturn ONLY the corrected JSON, no commentary, no markdown fences, in\nexactly this structure:\n\n{{"questions": [{{"question": "...", "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}}, "answer": "A", "explanation": "..."}}]}}\n\nJSON TO PROOFREAD:\n{json.dumps(payload, ensure_ascii=False)}\n'

def _proofread_batch(batch: list[dict], api_key: str, model: str, language: str) -> list[dict]:
    try:
        raw_reply = call_gemini(_build_proofread_prompt(batch), api_key, model)
        parsed = _parse_json_reply(raw_reply)
        fixed_list = parsed.get('questions') if isinstance(parsed, dict) else parsed
        if not isinstance(fixed_list, list) or len(fixed_list) != len(batch): return batch
    except Exception:
        return batch
    result = []
    for original, fixed in zip(batch, fixed_list):
        if not isinstance(fixed, dict) or not _is_valid_question(fixed):
            result.append(original)
            continue
        if fixed.get('answer') != original.get('answer'):
            result.append(original)
            continue
        candidate = dict(original)
        candidate['question'] = clean_math_formatting(clean_extraction_artifacts(fixed.get('question', original.get('question', ''))))
        candidate['options'] = {letter: clean_math_formatting(clean_extraction_artifacts(fixed.get('options', {}).get(letter, original.get('options', {}).get(letter, '')))) for letter in ('A', 'B', 'C', 'D')}
        candidate['explanation'] = clean_math_formatting(clean_extraction_artifacts(fixed.get('explanation', original.get('explanation', ''))))
        if _question_has_script_corruption(candidate, language):
            result.append(original)
        else:
            result.append(candidate)
    return result

async def _proofread_batch_limited(semaphore: asyncio.Semaphore, batch: list[dict], api_key: str, model: str, language: str) -> list[dict]:
    async with semaphore:
        return await asyncio.to_thread(_proofread_batch, batch, api_key, model, language)

async def proofread_questions(questions: list[dict], api_key: str, model: str, language: str) -> list[dict]:
    if not questions: return questions
    batch_size = 15
    sem = asyncio.Semaphore(GEMINI_MAX_CONCURRENCY)
    tasks = []
    for i in range(0, len(questions), batch_size):
        batch = questions[i:i + batch_size]
        tasks.append(_proofread_batch_limited(sem, batch, api_key, model, language))
    batches = await asyncio.gather(*tasks)
    proofed = []
    for b in batches:
        proofed.extend(b)
    return proofed

_VERIFY_CHUNK_SIZE = 20

def _build_verification_prompt(questions: list[dict], language: str) -> str:
    payload = [{'question': q.get('question', ''), 'options': q.get('options', {}), 'answer': q.get('answer', ''), 'explanation': q.get('explanation', '')} for q in questions]
    if language == 'Tamil':
        pattern_rules = '''
IMPORTANT ENFORCEMENT RULES FOR TAMIL EXAM PATTERNS:
1. For Statement-based questions (having statements 1., 2., 3. etc. in the question text):
   - The options A, B, C, D must be phrased using the exact Tamil convention:
     "1 மட்டும் சரி", "2 மட்டும் சரி", "1 மற்றும் 2 மட்டும் சரி", "2 மற்றும் 3 மட்டும் சரி", "1, 2 மற்றும் 3 சரி", "எதுவும் சரியில்லை" (or similar numbers matching the statements).
     NEVER use English words like "only" or "Both" or "None".
2. For Assertion-Reason questions (having "கூற்று (A):" and "காரணம் (R):" in the question text):
   - The options A, B, C, D must be phrased using the exact Tamil convention:
     "(A) சரி, (R) தவறு"
     "(A) தவறு, (R) சரி"
     "(A) மற்றும் (R) இரண்டும் சரி; (R) என்பது (A)-விற்கான சரியான விளக்கமாகும்"
     "(A) மற்றும் (R) இரண்டும் சரி; ஆனால் (R) என்பது (A)-வுக்கான சரியான விளக்கம் அல்ல"
     NEVER deviate from this exact wording or include English.
3. Clean any spelling or script errors (e.g. "உலோக்க" to "உலோகக்", "இலக்்க" to "இலக்க", "ப பங்கீட்டுப் பண்பு" to "பங்கீட்டுப் பண்பு").
'''
    else:
        pattern_rules = '''
IMPORTANT ENFORCEMENT RULES FOR ENGLISH EXAM PATTERNS:
1. For Statement-based questions (having statements Statement I, Statement II in the question text):
   - Options must use the standard conventions: "1 only", "2 only", "Both 1 and 2", "Neither 1 nor 2".
2. For Assertion-Reason questions (having "Assertion (A):" and "Reason (R):" in the question text):
   - Options must be:
     "A is true, R is false"
     "A is false, R is true"
     "Both A and R are true, and R is the correct explanation of A"
     "Both A and R are true, but R is not the correct explanation of A"
3. Clean any spelling errors and formatting issues.
'''
    return f'''You are a meticulous exam answer-key auditor working in {language}. Below is a JSON array of multiple-choice questions.

For EACH question, independently:
1. Solve it yourself from scratch using only the question and its four options.
2. Verify the correct answer option and check if the explanation is accurate, clear, and teaches the fact directly. The question, every option, and the explanation must NEVER refer to "the notes/passage/text/document/source/chapter/lesson/page/above" or any equivalent phrase in any language (e.g. Tamil "இந்த பக்கத்தில்", "குறிப்புகளின்படி") - if you find such a phrase, rewrite that field as a standalone fact/question with the reference removed entirely.
3. Rewrite the question, options, answer, and explanation if there are deviations from the rules.
4. Enforce the correct exam option patterns listed below.
5. If this is a mathematics question, enforce clean math formatting: never wrap
   variables/expressions in dollar signs ($x$, $$x$$) or LaTeX backslash commands
   (\\alpha, \\times, \\cdot) - use plain Unicode symbols instead (+ - × / = < > ≤ ≥ √ ^,
   and Greek letters α β γ θ λ μ π φ σ ω δ ε η ψ τ χ ξ ζ, never spelled out as words
   like "alpha" or "pi"), and vertical bars for absolute value (|x|). Never spell out
   operations in words ("divided by", "minus", "equals") - use the symbol. Silently
   correct any of these if you find them while rewriting.

{pattern_rules}

Return ONLY valid JSON, in exactly this structure:
{{
  "questions": [
    {{
      "question": "...",
      "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
      "answer": "A",
      "explanation": "..."
    }}
  ]
}}

QUESTIONS TO AUDIT:
{json.dumps(payload, ensure_ascii=False)}'''

def _verify_batch(batch: list[dict], api_key: str, model: str, language: str) -> list[dict]:
    try:
        raw_reply = call_gemini(_build_verification_prompt(batch, language), api_key, model)
        parsed = _parse_json_reply(raw_reply)
        fixed_list = parsed.get('questions') if isinstance(parsed, dict) else parsed
        if not isinstance(fixed_list, list) or len(fixed_list) != len(batch): return batch
    except Exception:
        return batch
    result = []
    for original, fixed in zip(batch, fixed_list):
        if not isinstance(fixed, dict) or not _is_valid_question(fixed):
            result.append(original)
            continue
        candidate = dict(original)
        candidate['question'] = clean_math_formatting(clean_extraction_artifacts(fixed.get('question', original.get('question', ''))))
        candidate['options'] = {letter: clean_math_formatting(clean_extraction_artifacts(fixed.get('options', {}).get(letter, original.get('options', {}).get(letter, '')))) for letter in ('A', 'B', 'C', 'D')}
        candidate['answer'] = fixed.get('answer', original.get('answer', ''))
        candidate['explanation'] = clean_math_formatting(clean_extraction_artifacts(fixed.get('explanation', original.get('explanation', ''))))
        if _question_has_script_corruption(candidate, language) or _question_references_source(candidate, language):
            result.append(original)
        else:
            result.append(candidate)
    return result

async def _verify_batch_limited(semaphore: asyncio.Semaphore, batch: list[dict], api_key: str, model: str, language: str) -> list[dict]:
    async with semaphore:
        return await asyncio.to_thread(_verify_batch, batch, api_key, model, language)

async def verify_and_correct_answers(questions: list[dict], api_key: str, model: str, language: str) -> list[dict]:
    if not questions: return questions
    batch_size = 15
    sem = asyncio.Semaphore(GEMINI_MAX_CONCURRENCY)
    tasks = []
    for i in range(0, len(questions), batch_size):
        batch = questions[i:i + batch_size]
        tasks.append(_verify_batch_limited(sem, batch, api_key, model, language))
    batches = await asyncio.gather(*tasks)
    verified = []
    for b in batches:
        verified.extend(b)
    return verified

async def get_notes_text(content: bytes, filename: str, ocr_lang: str='eng+tam'):
    file_hash = hashlib.sha256(content).hexdigest()
    cache_key = f'{file_hash}:{ocr_lang}'
    with _notes_text_cache_lock:
        cached = _notes_text_cache.get(cache_key)
    if cached is not None: return cached
    temp_dir = tempfile.gettempdir()
    temp_path = os.path.join(temp_dir, f'temp_{file_hash}_{filename}')
    async with aiofiles.open(temp_path, 'wb') as buffer:
        await buffer.write(content)
    try:
        notes_text = await asyncio.to_thread(extract_text, temp_path, ocr_lang)
        notes_text = clean_extraction_artifacts(notes_text)
        with _notes_text_cache_lock:
            _notes_text_cache[cache_key] = notes_text
        return notes_text
    finally:
        if os.path.exists(temp_path): os.remove(temp_path)

def split_text_into_windows(full_text: str, count: int, window_size: int=MAX_NOTES_CHARS) -> list[str]:
    n_chars = len(full_text)
    if count > 0 and n_chars > window_size * count:
        window_size = (n_chars // count) + 1
    if n_chars <= window_size: return [full_text]
    return [full_text[i:i + window_size] for i in range(0, n_chars, window_size)]

def split_count_across_windows(total_count: int, num_windows: int) -> list[int]:
    base = total_count // num_windows
    extra = total_count % num_windows
    return [base + 1 if i < extra else base for i in range(num_windows)]

def _split_count(count: int, chunk_size: int=CHUNK_SIZE) -> list[int]:
    chunks = []
    remaining = count
    while remaining > 0:
        take = min(chunk_size, remaining)
        chunks.append(take)
        remaining -= take
    return chunks

def _parse_json_reply(raw_reply: str):
    clean = raw_reply.strip()
    if clean.startswith('```'):
        clean = clean.strip('`')
        clean = clean.replace('json', '', 1).strip()
    return json.loads(clean)

_NEAR_DUPLICATE_THRESHOLD = 0.8
_STOPWORDS = {'the', 'a', 'an', 'is', 'are', 'was', 'were', 'of', 'to', 'in', 'on', 'and', 'or', 'for', 'with', 'by', 'which', 'what', 'who', 'whom', 'according', 'notes', 'did', 'does', 'do', 'as', 'that', 'this', 'these', 'those', 'it', 'at', 'from', 'his', 'her', 'their', 'its', 'type', 'kind', 'name', 'given', 'term', 'following', 'considered', 'described', 'provide', 'provides', 'very', 'little', 'primarily', 'assertion', 'reason', 'true', 'false', 'correct', 'explanation', 'statement', 'statements', 'neither', 'both', 'only'}
_TAMIL_STOPWORDS = {'என்ன', 'எது', 'எவை', 'எவ்வாறு', 'எத்தனை', 'யாது', 'ஆகும்', 'உள்ளது', 'உள்ளன', 'கிடைக்கும்', 'போது', 'ஒரு', 'ஓர்', 'மற்றும்', 'இந்த', 'அந்த', 'இவை', 'அவை', 'என்பது', 'என்று', 'என', 'எனப்படும்', 'எனப்பட', 'செய்ய', 'செய்தால்', 'கொண்டு', 'பயன்படுகிறது', 'பயன்படும்', 'குறிக்கிறது', 'குறிக்கும்', 'அழைக்கப்படுகிறது', 'கூறு', 'கூறுக', 'பின்வருவனவற்றுள்', 'பின்வரும்'}
_STOPWORDS |= _TAMIL_STOPWORDS
_WORD_RE = re.compile('[a-zA-Z0-9\\u0B80-\\u0BFF]+')
_SEMANTIC_OVERLAP_THRESHOLD = 0.4
_ANSWER_TEXT_DUPLICATE_THRESHOLD = 0.72

def _significant_words(text: str) -> set:
    text = (text or '').replace('-', '')
    words = _WORD_RE.findall(text.lower())
    return {w for w in words if len(w) >= 3 and w not in _STOPWORDS}

def _correct_answer_text(q: dict) -> str:
    options = q.get('options') or {}
    answer_letter = q.get('answer')
    return normalize_question_key(options.get(answer_letter, '')) if answer_letter else ''

def _is_near_duplicate(candidate: dict, seen: list) -> bool:
    for existing in seen:
        ratio = difflib.SequenceMatcher(None, candidate['key'], existing['key']).ratio()
        if ratio >= _NEAR_DUPLICATE_THRESHOLD:
            return True
        if ratio >= 0.68:
            answer_a, answer_b = (candidate.get('answer_text', ''), existing.get('answer_text', ''))
            if answer_a and answer_b and answer_a == answer_b:
                return True
    return False

_REQUIRED_OPTION_LETTERS = {'A', 'B', 'C', 'D'}
_ALLOWED_EXTRA_CHARS = (
    "\u2018\u2019\u201C\u201D\u2013\u2014\u2026\u00A0\u200c\u200d"
    "\u00d7\u00f7\u2212\u00b1\u00b0\u2032\u2033\u221a\u03c0"
    "\u00bc\u00bd\u00be\u00b2\u00b3\u20b9"
)
_LANGUAGE_SCRIPT_RANGES = {'Tamil': [(0x0B80, 0x0BFF)]}

_LATIN_WORD_RE = re.compile(r'[A-Za-z]{2,}')
_LATIN_WORD_ALLOWLIST = {
    'CM', 'MM', 'KM', 'KG', 'GM', 'ML', 'HZ', 'DNA', 'RNA', 'PH', 'AC', 'DC',
    'AM', 'PM', 'ID',
}

def _has_stray_latin_words(text_value: str, language: str) -> bool:
    if language != 'Tamil' or not text_value:
        return False
    for word in _LATIN_WORD_RE.findall(text_value):
        if word.upper() in _LATIN_WORD_ALLOWLIST:
            continue
        return True
    return False

def _question_has_script_corruption(q: dict, language: str) -> bool:
    fields = [q.get('question', ''), q.get('explanation', '')]
    fields.extend((q.get('options') or {}).values())
    for field in fields:
        if has_wrong_script_chars(field, language):
            return True
        if _has_stray_latin_words(field, language):
            return True
    return False

_SOURCE_REFERENCE_PATTERNS_EN = [
    re.compile(r'\baccording to the (notes|passage|text|document|material|source|above|given text)\b', re.IGNORECASE),
    re.compile(r'\bas (mentioned|stated|given|discussed|noted) (in|above)\b', re.IGNORECASE),
    re.compile(r'\bas per the (notes|passage|text|document|material|source|above)\b', re.IGNORECASE),
    re.compile(r'\bfrom the (notes|passage|text|document|material|source|above)\b', re.IGNORECASE),
    re.compile(r'\bin the (notes|passage|given text|document|material|source)\b', re.IGNORECASE),
    re.compile(r'\bbased on the (notes|passage|text|document|material|source|above)\b', re.IGNORECASE),
    re.compile(r'\bthe (notes|passage|text|document|material|source)\s+(state|states|say|says|mention|mentions|show|shows)\b', re.IGNORECASE),
    re.compile(r'\bin this (page|passage|note|notes|lesson|section|chapter|unit|paragraph)\b', re.IGNORECASE),
    re.compile(r'\bthis (page|passage|paragraph)\s+(states|mentions|says)\b', re.IGNORECASE),
    re.compile(r'\b(above|given)\s+(notes|passage|text|paragraph)\b', re.IGNORECASE),
]
_SOURCE_REFERENCE_SUBSTRINGS_TA = [
    'இந்த பக்கத்தில்', 'பக்கத்தில் இருந்து', 'இந்த குறிப்பில்',
    'குறிப்புகளின்படி', 'குறிப்புகளின் படி', 'இந்த பாடத்தில்',
    'மேலே கொடுக்கப்பட்ட', 'கொடுக்கப்பட்ட பத்தியில்', 'இந்த பத்தியில்',
]

def _question_references_source(q: dict, language: str) -> bool:
    fields = [q.get('question', ''), q.get('explanation', '')]
    fields.extend((q.get('options') or {}).values())
    combined = ' '.join((f for f in fields if f))
    if not combined:
        return False
    if any(pattern.search(combined) for pattern in _SOURCE_REFERENCE_PATTERNS_EN):
        return True
    if any(substr in combined for substr in _SOURCE_REFERENCE_SUBSTRINGS_TA):
        return True
    return False

def _is_valid_question(q: dict) -> bool:
    if not isinstance(q, dict): return False
    if not q.get('question'): return False
    options = q.get('options')
    if not isinstance(options, dict) or not _REQUIRED_OPTION_LETTERS.issubset(options.keys()): return False
    if not all((options.get(letter) for letter in _REQUIRED_OPTION_LETTERS)): return False
    if q.get('answer') not in _REQUIRED_OPTION_LETTERS: return False
    return True

_MAX_BACKFILL_ATTEMPTS = 10
_BACKFILL_OVERASK_FACTOR = 2.5
_BACKFILL_CHUNK_CAP = 20
_MAX_AVOID_LIST_LEN = 60

def _process_single_question(q: dict, language: str, seen: list) -> str:
    if not _is_valid_question(q) or _question_has_script_corruption(q, language):
        return "malformed"
    if _question_references_source(q, language):
        return "malformed"
    q['question'] = clean_math_formatting(clean_extraction_artifacts(q.get('question', '')))
    q['options'] = {letter: clean_math_formatting(clean_extraction_artifacts(q['options'][letter])) for letter in ('A', 'B', 'C', 'D')}
    if (q.get('explanation') or '').strip():
        q['explanation'] = clean_math_formatting(clean_extraction_artifacts(q['explanation']))
    else:
        q['explanation'] = f"The correct answer is {q['options'][q['answer']]}."
    key = normalize_question_key(q.get('question', ''))
    if not key:
        return "malformed"
    answer_text = _correct_answer_text(q)
    answer_words = _significant_words(answer_text)
    candidate = {'key': key, 'words': _significant_words(q.get('question', '')) | answer_words, 'answer_words': answer_words, 'answer_text': answer_text}
    if _is_near_duplicate(candidate, seen):
        return "duplicate"
    seen.append(candidate)
    return "ok"

def _merge_chunk_replies(raw_replies: list[str], seen: list, merged_questions: list, language: str) -> tuple[int, int]:
    dropped_malformed = 0
    dropped_duplicate = 0
    for raw_reply in raw_replies:
        try:
            chunk_data = _parse_json_reply(raw_reply)
        except ValueError:
            continue
        if isinstance(chunk_data, dict):
            questions_list = chunk_data.get('questions', [])
        elif isinstance(chunk_data, list):
            questions_list = chunk_data
        else:
            questions_list = []
        for q in questions_list:
            status = _process_single_question(q, language, seen)
            if status == "ok":
                merged_questions.append(q)
            elif status == "duplicate":
                dropped_duplicate += 1
            else:
                dropped_malformed += 1
    return (dropped_malformed, dropped_duplicate)

async def _call_gemini_limited(semaphore: asyncio.Semaphore, prompt: str, api_key: str, model: str):
    async with semaphore:
        return await asyncio.to_thread(call_gemini, prompt, api_key, model)

_INITIAL_OVERASK_FACTOR = 1.2

async def generate_questions_data(notes_text: str, count: int, difficulty: str, api_key: str, model: str=DEFAULT_GEMINI_MODEL, language: str=None) -> dict:
    initial_ask = max(count, round(count * _INITIAL_OVERASK_FACTOR))
    chunk_size = GENERATION_CHUNK_SIZE
    chunk_threshold = chunk_size + 2
    chunk_sizes = _split_count(initial_ask, chunk_size) if initial_ask > chunk_threshold else [initial_ask]
    if language is None:
        language = detect_dominant_language(notes_text)
    merged_questions: list[dict] = []
    seen: list[dict] = []
    sem = asyncio.Semaphore(GEMINI_MAX_CONCURRENCY)
    prompts = []
    current_fmt_idx = 0
    for idx, size in enumerate(chunk_sizes):
        formats_list = []
        for _ in range(size):
            formats_list.append((current_fmt_idx % 14) + 1)
            current_fmt_idx += 1
        diversity_hint = ""
        if len(chunk_sizes) > 1:
            diversity_hint = f"\n\nADDITIONAL CONTEXT FOR DIVERSITY:\nThis is request {idx + 1} of {len(chunk_sizes)} concurrent generation requests for these notes.\nTo avoid duplicates, focus on a different part or aspect of the notes (e.g., section {idx + 1} of the content) and write unique questions.\n"
        prompt_str = build_prompt(notes_text, size, difficulty, language, formats_list) + diversity_hint
        prompts.append(prompt_str)
    tasks = [_call_gemini_limited(sem, p, api_key, model) for p in prompts]
    raw_replies = await asyncio.gather(*tasks)
    await asyncio.to_thread(_merge_chunk_replies, raw_replies, seen, merged_questions, language)
    backfill_attempts = 0
    while len(merged_questions) < count and backfill_attempts < 5:
        shortfall = count - len(merged_questions)
        ask_for = min(25, max(shortfall, round(shortfall * 2.0)))
        backfill_formats = []
        for i in range(ask_for):
            backfill_formats.append(((len(merged_questions) + i) % 14) + 1)
        avoid_list = [q['question'] for q in merged_questions[-60:]]
        backfill_prompt = build_prompt(notes_text, ask_for, difficulty, language, backfill_formats, avoid_questions=avoid_list)
        try:
            backfill_reply = await asyncio.to_thread(call_gemini, backfill_prompt, api_key, model)
            _merge_chunk_replies([backfill_reply], seen, merged_questions, language)
        except Exception:
            pass
        backfill_attempts += 1
    delivered = merged_questions[:count]
    return {'questions': delivered, 'requested_count': count, 'delivered_count': len(delivered)}

_ANSWER_LETTER_TO_NUM = {'A': '1', 'B': '2', 'C': '3', 'D': '4'}
SUBJECT_QUES_PREFIX_MAP = {'Tamil': 'TA', 'English': 'EN', 'Maths': 'MA', 'Science': 'SC', 'Social Science': 'SO', 'Physics': 'PH', 'Chemistry': 'CH', 'Biology': 'BI', 'Computer Science': 'CS', 'Botany': 'BO', 'Zoology': 'ZO', 'Commerce': 'CO', 'Economics': 'EC', 'Accountancy': 'AC', 'Computer Application': 'CA', 'Mathematics': 'MA'}

def _build_dy_rows(questions: list[dict], subject: str, standard: str, dy_code: str, ques_id_prefix: str) -> list[dict]:
    # dy_code and ques_id_prefix are two separate, independently-typed fields:
    #   - dy_code       -> the exam code, saved as-is (used for the order/serial
    #                       counter below, and as the exam identifier elsewhere).
    #   - ques_id_prefix -> the literal prefix used to build dy_ques_id, e.g.
    #                       typing "AI8EN2" here makes IDs AI8EN2001, AI8EN2002...
    # They used to be silently merged into one auto-computed value (standard +
    # subject only), which ignored whatever the user actually typed and caused
    # different exams/weeks to generate colliding dy_ques_id values. Now both
    # are taken exactly as typed.
    dy_code = (dy_code or '').strip().upper()
    id_prefix = (ques_id_prefix or '').strip().upper()
    rows = []

    # Get sequential order numbers for this dy_code (1, 2, 3...)
    order_numbers = get_next_order_numbers_batch(dy_code, len(questions))

    for q, order_number in zip(questions, order_numbers):
        options = q.get('options') or {}
        rows.append({
            'dy_ques_id': f'{id_prefix}{str(order_number).zfill(3)}',
            'dy_code': dy_code,
            'ln_code': 'LN01' if _TAMIL_CHAR_RE.search(q.get('question', '')) else 'LN02',
            'dy_order': str(order_number),
            'dy_pattern': '1',
            'dy_seconds': '18',
            'dy_question': q.get('question', ''),
            'dy_image_name': None,
            'dy_ans_1': options.get('A', ''),
            'dy_ans_2': options.get('B', ''),
            'dy_ans_3': options.get('C', ''),
            'dy_ans_4': options.get('D', ''),
            'dy_correct_ans': _ANSWER_LETTER_TO_NUM.get(q.get('answer', ''), ''),
            'dy_explain': q.get('explanation', ''),
            'dy_explain_image_name': None
        })
    return rows

def _get_selected_ids(db, question_ids: list[int] | None=None) -> set[int]:
    query = db.query(SelectedQuestion.question_id)
    if question_ids is not None: query = query.filter(SelectedQuestion.question_id.in_(question_ids))
    return {row[0] for row in query.all()}

def _serialize_question(record: QuestionRecord, selected_ids: set[int]) -> dict:
    try:
        pattern_val = int(record.dy_pattern) if record.dy_pattern and record.dy_pattern.strip().isdigit() else 1
    except Exception:
        pattern_val = 1
    try:
        seconds_val = int(record.dy_seconds) if record.dy_seconds and record.dy_seconds.strip().isdigit() else 18
    except Exception:
        seconds_val = 18
    return {
        'id': record.id,
        'source_filename': record.source_filename,
        'dy_ques_id': record.dy_ques_id,
        'dy_code': record.dy_code,
        'ln_code': record.ln_code,
        'dy_order': record.dy_order,
        'dy_pattern': pattern_val,
        'dy_seconds': seconds_val,
        'dy_question': record.dy_question,
        'dy_image_name': record.dy_image_name or "",
        'dy_ans_1': record.dy_ans_1,
        'dy_ans_2': record.dy_ans_2,
        'dy_ans_3': record.dy_ans_3,
        'dy_ans_4': record.dy_ans_4,
        'dy_correct_ans': record.dy_correct_ans,
        'dy_explain': record.dy_explain,
        'dy_explain_image_name': record.dy_explain_image_name or "",
        'created_at': record.created_at.isoformat() if record.created_at else None,
        'updated_at': record.updated_at.isoformat() if record.updated_at else None,
        'difficulty': record.difficulty,
        'batch_id': record.batch_id,
        'board': record.board,
        'standard': record.standard,
        'group_name': record.group_name,
        'subject': record.subject,
        'selected': record.id in selected_ids
    }

class SelectionPayload(BaseModel):
    batch_id: str
    question_ids: list[int]

class QuestionIdsPayload(BaseModel):
    question_ids: list[int]

class TagPayload(BaseModel):
    question_ids: list[int]
    board: str | None = None
    standard: str | None = None
    group_name: str | None = None
    subject: str | None = None

class PushToLiveRequest(BaseModel):
    exam_type: Literal['daily', 'schedule', 'online']
    batch_id: str | None = None
    dy_ques_ids: list[str] | None = None
    only_selected: bool = False
    exam_code: str | None = None
    force: bool = False

_LIVE_PUSH_JSON_KEYS = [
    'dy_ques_id', 'dy_code', 'ln_code', 'dy_order', 'dy_pattern', 'dy_seconds',
    'dy_question', 'dy_ans_1', 'dy_ans_2', 'dy_ans_3', 'dy_ans_4',
    'dy_correct_ans', 'dy_explain',
]

class QuestionEditPayload(BaseModel):
    question: str
    options: dict[str, str]
    answer: str
    explanation: str

def _expected_ocr_lang(subject: str, board: str) -> str:
    return 'eng+tam'

def _detect_generation_language(subject: str, board: str, full_notes_text: str) -> str:
    tamil_count = len(TAMIL_CHAR_RE.findall(full_notes_text))
    latin_count = len(LATIN_CHAR_RE.findall(full_notes_text))
    result = detect_dominant_language(full_notes_text)
    print(f"[lang-detect] subject={subject!r} board={board!r} (selection ignored for language) -> "
          f"tamil_chars={tamil_count} latin_chars={latin_count} notes_len={len(full_notes_text)} -> result={result!r}")
    return result

def _merge_window_results(window_results: list, seen: list, all_questions: list) -> None:
    for window_data in window_results:
        window_data = clean_question_payload(window_data)
        for q in window_data.get('questions', []):
            key = normalize_question_key(q.get('question', ''))
            answer_text = _correct_answer_text(q)
            answer_words = _significant_words(answer_text)
            candidate = {'key': key, 'words': _significant_words(q.get('question', '')) | answer_words, 'answer_words': answer_words, 'answer_text': answer_text}
            if _is_near_duplicate(candidate, seen):
                continue
            seen.append(candidate)
            all_questions.append(q)

async def _run_backfill(count: int, all_questions: list, seen: list, windows: list, difficulty: str, overall_language: str, api_key: str, selected_model: str) -> None:
    backfill_attempts = 0
    while len(all_questions) < count and backfill_attempts < 12:
        shortfall = count - len(all_questions)
        ask_for = min(25, max(shortfall, round(shortfall * 2.0)))
        window_idx = backfill_attempts % len(windows)
        target_window_text = windows[window_idx]
        avoid_list = [q['question'] for q in all_questions[-60:]]
        backfill_formats = []
        for i in range(ask_for):
            backfill_formats.append(((len(all_questions) + i) % 14) + 1)
        backfill_prompt = build_prompt(target_window_text, ask_for, difficulty, overall_language, backfill_formats, avoid_questions=avoid_list)
        try:
            raw_reply = await asyncio.to_thread(call_gemini, backfill_prompt, api_key, selected_model)
            chunk_data = _parse_json_reply(raw_reply)
        except Exception:
            backfill_attempts += 1
            continue
        questions_list = chunk_data.get('questions', []) if isinstance(chunk_data, dict) else chunk_data
        if not isinstance(questions_list, list):
            backfill_attempts += 1
            continue
        for q in questions_list:
            if len(all_questions) >= count:
                break
            status = _process_single_question(q, overall_language, seen)
            if status == "ok":
                all_questions.append(q)
        backfill_attempts += 1

def _split_out_db_duplicates(questions: list[dict], filename: str) -> tuple[list[dict], int]:
    db = SessionLocal()
    try:
        query = db.query(QuestionRecord)
        if filename:
            query = query.filter(QuestionRecord.source_filename == filename)
        else:
            query = query.filter(QuestionRecord.source_filename.is_(None))
        existing_records = query.all()
    finally:
        db.close()

    letter_map = {'1': 'A', '2': 'B', '3': 'C', '4': 'D'}
    precomputed_existing = []
    for r in existing_records:
        ekey = normalize_question_key(r.dy_question)
        correct_letter = letter_map.get(r.dy_correct_ans, '')
        opts = {'A': r.dy_ans_1, 'B': r.dy_ans_2, 'C': r.dy_ans_3, 'D': r.dy_ans_4}
        correct_txt = normalize_question_key(opts.get(correct_letter, '')) if correct_letter else ''
        precomputed_existing.append({'key': ekey, 'answer_text': correct_txt})

    kept = []
    skipped_existing = 0
    for q in questions:
        key = normalize_question_key(q.get('question', ''))
        answer_text = _correct_answer_text(q)

        is_dup = False
        for existing in precomputed_existing:
            ratio = difflib.SequenceMatcher(None, key, existing['key']).ratio()
            if ratio >= _NEAR_DUPLICATE_THRESHOLD:
                is_dup = True
                break
            if ratio >= 0.68:
                existing_answer = existing.get('answer_text', '')
                if answer_text and existing_answer and answer_text == existing_answer:
                    is_dup = True
                    break

        if is_dup:
            skipped_existing += 1
            q['id'] = None
            q['dy_ques_id'] = None
            q['duplicate'] = True
        else:
            kept.append(q)
    return kept, skipped_existing

def _save_generated_questions(questions: list[dict], rows: list[dict], difficulty: str, filename: str, batch_id: str, selected_model: str, board: str, standard: str, subject: str, group_name: str | None) -> None:
    db = SessionLocal()
    try:
        for q, row in zip(questions, rows):
            record = QuestionRecord(dy_ques_id=row['dy_ques_id'], dy_code=row['dy_code'], ln_code=row['ln_code'], dy_order=row['dy_order'], dy_pattern=row['dy_pattern'], dy_seconds=row['dy_seconds'], dy_question=row['dy_question'], dy_image_name=row['dy_image_name'], dy_ans_1=row['dy_ans_1'], dy_ans_2=row['dy_ans_2'], dy_ans_3=row['dy_ans_3'], dy_ans_4=row['dy_ans_4'], dy_correct_ans=row['dy_correct_ans'], dy_explain=row['dy_explain'], dy_explain_image_name=row['dy_explain_image_name'], difficulty=difficulty, source_filename=filename, batch_id=batch_id, model=selected_model, board=board, standard=standard, subject=subject, group_name=group_name)
            db.add(record)
            db.flush()
            q['id'] = record.id
            q['dy_ques_id'] = row['dy_ques_id']
            q['dy_code'] = row['dy_code']
            q['ln_code'] = row['ln_code']
            q['dy_order'] = row['dy_order']
            q['dy_pattern'] = row['dy_pattern']
            q['dy_seconds'] = row['dy_seconds']
            q['duplicate'] = False
        db.commit()
    finally:
        db.close()

@app.post('/generate-questions', responses={400: {'description': 'Invalid difficulty, count, or model.'}, 500: {'description': 'GEMINI_API_KEY is not configured on the server.'}, 502: {'description': 'Gemini API error.'}})
async def generate_questions_endpoint(file: Annotated[UploadFile, File()], count: Annotated[int, Form()], difficulty: Annotated[str, Form()], subject: Annotated[str, Form()], board: Annotated[str, Form()], standard: Annotated[str, Form()], dy_code: Annotated[str, Form()], ques_id_prefix: Annotated[str, Form()], group_name: Annotated[str | None, Form()]=None, model: Annotated[str | None, Form()]=None):
    difficulty = difficulty.lower().strip()
    if difficulty not in ('easy', 'moderate', 'hard'): raise HTTPException(status_code=400, detail='difficulty must be easy, moderate, or hard')
    if count < 1 or count > 500: raise HTTPException(status_code=400, detail='count must be between 1 and 500')
    selected_model = (model or DEFAULT_GEMINI_MODEL).strip()
    if selected_model not in ALLOWED_GEMINI_MODELS: raise HTTPException(status_code=400, detail=f'model must be one of: {", ".join(ALLOWED_GEMINI_MODELS)}')
    subject = (subject or '').strip()
    board = (board or '').strip()
    standard = (standard or '').strip()
    dy_code = (dy_code or '').strip().upper()
    ques_id_prefix = (ques_id_prefix or '').strip().upper()
    group_name = (group_name or '').strip() or None
    if not subject or not board or (not standard): raise HTTPException(status_code=400, detail='board, standard and subject are required')
    if not dy_code: raise HTTPException(status_code=400, detail='dy_code is required')
    if not ques_id_prefix: raise HTTPException(status_code=400, detail='ques_id_prefix is required')
    content = await file.read()
    ocr_lang = _expected_ocr_lang(subject, board)
    print(f"[ocr] subject={subject!r} board={board!r} -> ocr_lang={ocr_lang!r}")
    full_notes_text = await get_notes_text(content, file.filename, ocr_lang)
    if not full_notes_text.strip(): raise HTTPException(status_code=400, detail='No text could be extracted from this file')
    windows = split_text_into_windows(full_notes_text, count)
    window_counts = split_count_across_windows(count, len(windows))
    api_key = os.getenv('GEMINI_API_KEY')
    if not api_key: raise HTTPException(status_code=500, detail='GEMINI_API_KEY is not set on the server.')
    overall_language = _detect_generation_language(subject, board, full_notes_text)
    tasks = []
    for window_text, window_count in zip(windows, window_counts):
        if window_count <= 0: continue
        tasks.append(generate_questions_data(window_text, window_count, difficulty, api_key, selected_model, language=overall_language))
    try:
        window_results = await asyncio.gather(*tasks)
    except errors.ClientError as e:
        if getattr(e, 'code', None) == 404 and 'no longer available' in str(e).lower(): raise HTTPException(status_code=502, detail=f"Model '{selected_model}' (and the fallback model) are not available on this Gemini API key. Try a different model from GET /models.")
        raise HTTPException(status_code=502, detail=f'Gemini API error: {e}')
    all_questions: list[dict] = []
    seen: list[dict] = []
    _merge_window_results(window_results, seen, all_questions)
    await _run_backfill(count, all_questions, seen, windows, difficulty, overall_language, api_key, selected_model)
    all_questions = await verify_and_correct_answers(all_questions, api_key, selected_model, overall_language)
    all_questions = await proofread_questions(all_questions, api_key, selected_model, overall_language)
    data = {'questions': all_questions, 'requested_count': count, 'delivered_count': len(all_questions)}
    data = fix_answer_consistency(data)
    batch_id = str(uuid.uuid4())
    kept_questions, skipped_existing = _split_out_db_duplicates(data['questions'], file.filename)
    rows = _build_dy_rows(kept_questions, subject, standard, dy_code, ques_id_prefix)
    _save_generated_questions(kept_questions, rows, difficulty, file.filename, batch_id, selected_model, board, standard, subject, group_name)
    shortfall_note = None
    if data['delivered_count'] < data['requested_count']: shortfall_note = f'Only {data["delivered_count"]} of {data["requested_count"]} requested questions could be generated without repeating the same underlying fact.'
    duplicate_note = None
    if skipped_existing: duplicate_note = f'{skipped_existing} question(s) already existed in the database for this file and difficulty, so they were not saved again.'
    return {'questions': data['questions'], 'count': data['delivered_count'], 'requested_count': data['requested_count'], 'difficulty': difficulty, 'batch_id': batch_id, 'model': selected_model, 'note': shortfall_note, 'duplicate_note': duplicate_note, 'board': board, 'standard': standard, 'subject': subject, 'group_name': group_name, 'dy_code': dy_code}

@app.get('/models')
def list_models():
    return {'default': DEFAULT_GEMINI_MODEL, 'models': [{'id': model_id, 'label': label} for model_id, label in ALLOWED_GEMINI_MODELS.items()]}

@app.post('/questions/select', responses={400: {'description': 'batch_id is required'}})
def save_selection(payload: SelectionPayload):
    batch_id = (payload.batch_id or '').strip()
    if not batch_id: raise HTTPException(status_code=400, detail='batch_id is required')
    question_ids = list(dict.fromkeys(payload.question_ids))
    db = SessionLocal()
    try:
        valid_records = db.query(QuestionRecord).filter(QuestionRecord.id.in_(question_ids), QuestionRecord.batch_id == batch_id).all()
        db.query(SelectedQuestion).filter(SelectedQuestion.batch_id == batch_id).delete()
        db.add_all([SelectedQuestion(question_id=r.id, batch_id=batch_id) for r in valid_records])
        db.commit()
        return {'status': 'saved', 'batch_id': batch_id, 'selected_count': len(valid_records)}
    finally:
        db.close()

@app.delete('/questions/batch/{batch_id}', responses={404: {'description': 'No matching batch found'}})
def delete_batch(batch_id: str): 
    db = SessionLocal()
    try:
        db.query(SelectedQuestion).filter(SelectedQuestion.batch_id == batch_id).delete()
        deleted_count = db.query(QuestionRecord).filter(QuestionRecord.batch_id == batch_id).delete()
        db.commit()
        if deleted_count == 0: raise HTTPException(status_code=404, detail='No matching batch found.')
        return {'status': 'deleted', 'batch_id': batch_id, 'deleted_count': deleted_count}
    finally:
        db.close()

@app.post('/questions/deselect', responses={400: {'description': _ERR_EMPTY_QUESTION_IDS}})
def deselect_questions(payload: QuestionIdsPayload):
    question_ids = list(dict.fromkeys(payload.question_ids))
    if not question_ids: raise HTTPException(status_code=400, detail=_ERR_EMPTY_QUESTION_IDS)
    db = SessionLocal()
    try:
        deselected_count = db.query(SelectedQuestion).filter(SelectedQuestion.question_id.in_(question_ids)).delete(synchronize_session=False)
        db.commit()
        return {'status': 'deselected', 'deselected_count': deselected_count}
    finally:
        db.close()

@app.put('/questions/{question_id}', responses={400: {'description': 'Invalid question text, options, or answer'}, 404: {'description': 'Question not found'}})
def edit_question(question_id: int, payload: QuestionEditPayload):
    question_text = (payload.question or '').strip()
    answer = (payload.answer or '').strip().upper()
    options = payload.options or {}
    if not question_text: raise HTTPException(status_code=400, detail='Question text cannot be empty')
    if answer not in _ANSWER_LETTER_TO_NUM: raise HTTPException(status_code=400, detail='answer must be A, B, C, or D')
    if not all(((options.get(letter) or '').strip() for letter in ('A', 'B', 'C', 'D'))): raise HTTPException(status_code=400, detail='All four options are required')
    db = SessionLocal()
    try:
        record = db.query(QuestionRecord).filter(QuestionRecord.id == question_id).first()
        if not record: raise HTTPException(status_code=404, detail='Question not found')
        record.dy_question = question_text
        record.dy_ans_1 = options['A'].strip()
        record.dy_ans_2 = options['B'].strip()
        record.dy_ans_3 = options['C'].strip()
        record.dy_ans_4 = options['D'].strip()
        record.dy_correct_ans = _ANSWER_LETTER_TO_NUM[answer]
        record.dy_explain = (payload.explanation or '').strip()
        record.updated_at = datetime.now(timezone.utc).replace(tzinfo=None)
        db.commit()
        return {'status': 'updated', 'id': record.id}
    finally:
        db.close()

@app.delete('/questions/{question_id}', responses={404: {'description': 'Question not found'}})
def delete_question(question_id: int):
    db = SessionLocal()
    try:
        record = db.query(QuestionRecord).filter(QuestionRecord.id == question_id).first()
        if not record: raise HTTPException(status_code=404, detail='Question not found')
        db.query(SelectedQuestion).filter(SelectedQuestion.question_id == question_id).delete(synchronize_session=False)
        db.delete(record)
        db.commit()
        return {'status': 'deleted', 'id': question_id}
    finally:
        db.close()

@app.post('/questions/tag', responses={400: {'description': _ERR_EMPTY_QUESTION_IDS}})
def tag_questions(payload: TagPayload):
    question_ids = list(dict.fromkeys(payload.question_ids))
    if not question_ids: raise HTTPException(status_code=400, detail=_ERR_EMPTY_QUESTION_IDS)
    db = SessionLocal()
    try:
        records = db.query(QuestionRecord).filter(QuestionRecord.id.in_(question_ids)).all()
        for record in records:
            if payload.board is not None: record.board = payload.board
            if payload.standard is not None: record.standard = payload.standard
            if payload.group_name is not None: record.group_name = payload.group_name
            if payload.subject is not None: record.subject = payload.subject
        existing_selected = _get_selected_ids(db, [r.id for r in records])
        db.add_all([SelectedQuestion(question_id=r.id) for r in records if r.id not in existing_selected])
        db.commit()
        return {'status': 'tagged', 'count': len(records)}
    finally:
        db.close() 

@app.post('/questions/untag', responses={400: {'description': _ERR_EMPTY_QUESTION_IDS}})
def untag_questions(payload: QuestionIdsPayload):
    question_ids = list(dict.fromkeys(payload.question_ids))
    if not question_ids: raise HTTPException(status_code=400, detail=_ERR_EMPTY_QUESTION_IDS)
    db = SessionLocal()
    try:
        db.query(QuestionRecord).filter(QuestionRecord.id.in_(question_ids)).update({QuestionRecord.board: None, QuestionRecord.standard: None, QuestionRecord.group_name: None, QuestionRecord.subject: None}, synchronize_session=False)
        db.query(SelectedQuestion).filter(SelectedQuestion.question_id.in_(question_ids)).delete(synchronize_session=False)
        db.commit()
        return {'status': 'untagged', 'count': len(question_ids)}
    finally:
        db.close()

@app.get('/questions')
def get_all_questions():
    db = SessionLocal()
    try:
        records = db.query(QuestionRecord).order_by(QuestionRecord.created_at.desc()).all()
        selected_ids = _get_selected_ids(db, [r.id for r in records])
        return {'questions': [_serialize_question(r, selected_ids) for r in records]}
    finally:
        db.close()

@app.post('/push-to-live', responses={
    400: {'description': 'Provide either batch_id or dy_ques_ids.'},
    404: {'description': 'No matching questions found.'},
    500: {'description': 'EXTERNAL_DB_API_URL / EXTERNAL_API_KEY are not configured on the server.'},
})
async def push_to_live(payload: PushToLiveRequest):
    if not payload.batch_id and not payload.dy_ques_ids:
        raise HTTPException(status_code=400, detail='Provide either batch_id (push a whole batch) or dy_ques_ids (push a specific selection).')

    def load_records():
        db = SessionLocal()
        try:
            if payload.dy_ques_ids:
                return db.query(QuestionRecord).filter(QuestionRecord.dy_ques_id.in_(payload.dy_ques_ids)).order_by(cast(QuestionRecord.dy_order, Integer)).all()
            query = db.query(QuestionRecord).filter(QuestionRecord.batch_id == payload.batch_id)
            if payload.only_selected:
                selected_ids = {s.question_id for s in db.query(SelectedQuestion).filter(SelectedQuestion.batch_id == payload.batch_id).all()}
                query = query.filter(QuestionRecord.id.in_(selected_ids))
            return query.order_by(cast(QuestionRecord.dy_order, Integer)).all()
        finally:
            db.close()

    records = await asyncio.to_thread(load_records)
    if not records:
        raise HTTPException(status_code=404, detail='No matching questions found.')

    try:
        result = await push_records_to_live_db(records, payload.exam_type, payload.exam_code, payload.force)
    except RuntimeError as e:
        raise HTTPException(status_code=500, detail=str(e))

    return result

@app.post('/push-json-to-live', responses={
    400: {'description': 'Invalid or empty JSON file.'},
    500: {'description': 'EXTERNAL_DB_API_URL / EXTERNAL_API_KEY are not configured on the server.'},
})
async def push_json_to_live(
    file: Annotated[UploadFile, File()],
    exam_type: Annotated[Literal['daily', 'schedule', 'online'], Form()],
    exam_code: Annotated[str | None, Form()] = None,
    force: Annotated[bool, Form()] = False,
):
    raw_bytes = await file.read()
    try:
        data = json.loads(raw_bytes.decode('utf-8'))
    except (json.JSONDecodeError, UnicodeDecodeError) as e:
        raise HTTPException(status_code=400, detail=f'Could not parse JSON file: {e}')

    if isinstance(data, dict) and 'questions' in data:
        data = data['questions']
    if not isinstance(data, list) or not data:
        raise HTTPException(status_code=400, detail='Expected a non-empty JSON array of questions, or {"questions": [...]}.')

    records = []
    for i, item in enumerate(data):
        if not isinstance(item, dict):
            raise HTTPException(status_code=400, detail=f'Item {i} in the JSON file is not an object.')
        safe = {key: item.get(key) for key in _LIVE_PUSH_JSON_KEYS}
        records.append(SimpleNamespace(**safe))

    try:
        result = await push_records_to_live_db(records, exam_type, exam_code, force)
    except RuntimeError as e:
        raise HTTPException(status_code=500, detail=str(e))

    return result

@app.get('/live-push-status')
async def live_push_status(
    exam_type: Annotated[Literal['daily', 'schedule', 'online'], Query()],
    batch_id: Annotated[str | None, Query()] = None,
):
    def fetch():
        db = SessionLocal()
        try:
            query = db.query(LivePushLog).filter(LivePushLog.exam_type == exam_type)
            if batch_id:
                batch_ids = {r.dy_ques_id for r in db.query(QuestionRecord).filter(QuestionRecord.batch_id == batch_id).all()}
                query = query.filter(LivePushLog.dy_ques_id.in_(batch_ids))
            return query.all()
        finally:
            db.close()

    rows = await asyncio.to_thread(fetch)
    return {
        'status': 'success',
        'exam_type': exam_type,
        'records': [
            {'dy_ques_id': r.dy_ques_id, 'status': r.status, 'live_ques_id': r.live_ques_id, 'pushed_at': r.pushed_at}
            for r in rows
        ],
    }